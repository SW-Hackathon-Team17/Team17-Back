# from datetime import datetime
# from http import HTTPStatus
# from os import access
# from flask import request
# from flask_jwt_extended import create_access_token, get_jwt, get_jwt_identity, jwt_required
# from flask_restful import Resource
# from mysql.connector.errors import Error
# import boto3

# from config import Config

# class FileUploadResource(Resource) :

#     def post(self) :

#         # 1. 클라이언트로부터 데이터를 받아온다.
#         # request.files 에 파일을 받아온다.
#         # 따라서 파일이 없는 상태로 API 가 호출되면, 에러메세지를
#         # 클라이언트에 응답해주자.

#         # photo란, 클라이언트에서 보내는 key !!
#         if 'photo' not in request.files:
#             return {'error' : '파일을 업로드 하세요'}, 400

#         # 클라이언트로부터 파일을 받아온다.
#         file = request.files['photo']

#         # 파일명을 우리가 변경해 준다.
#         # 파일명은, 유니크하게 만들어야 한다.
#         current_time = datetime.now()
#         new_file_name = current_time.isoformat().replace(':', '_') + '.jpg'

#         # 유저가 올린 파일의 이름을, 내가 만든 파일명으로 변경
#         file.filename = new_file_name

#         # S3 에 업로드 하면 된다.
#         # AWS 의 라이브러리를 사용해야 한다.
#         # 이 파이썬 라이브러리가 boto3 라이브러리다!
#         # boto3 라이브러리 설치
#         # pip install boto3

#         s3 = boto3.client('s3',
#                     aws_access_key_id = Config.ACCESS_KEY,
#                     aws_secret_access_key = Config.SECRET_ACCESS)

#         try :
#             s3.upload_fileobj(file,
#                                 Config.S3_BUCKET,
#                                 file.filename,
#                                 ExtraArgs = {'ACL':'public-read', 'ContentType':file.content_type} )

#         except Exception as e :
#             return {'error' : str(e)}, 500

#         return {'result' : 'success',
#                 'imgUrl' : Config.S3_LOCATION + file.filename}


import os
from flask import Flask, Blueprint, request
from pptx import Presentation
from PIL import Image
import boto3
from io import BytesIO
import sqlite3

bp = Blueprint('image', __name__, url_prefix='/upload')

# app = Flask(__name__)

s3 = boto3.client('s3',
                  aws_access_key_id='AKIAW2M3VPW76QF656TM',
                  aws_secret_access_key='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')


def pptx_to_images(file):
    prs = Presentation(file)
    image_list = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                img = Image.open(BytesIO(image_bytes))
                image_list.append(img)
    return image_list


def upload_to_s3_and_save_to_db(images):
    # conn = sqlite3.connect('your_database.db')
    # cur = conn.cursor()
    li = []
    for i, image in enumerate(images):
        out_img = BytesIO()
        image.save(out_img, format='png')
        out_img.seek(0)

        key = f"image_{i}.png"

        s3.upload_fileobj(out_img, 'team17-buckets', key)

        url = f"https://team17-buckets.s3.amazonaws.com/{key}"
        li.append(url)

    return li
    # Assuming you have a table named 'images' with 'url' column

    # conn.commit()
    # cur.close()
    # conn.close()


@bp.route('/', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return 'No file part in the request'
    file = request.files['file']
    if file.filename == '':
        return 'No file selected for uploading'
    if file and file.filename.endswith('.pptx'):
        print(file)
        images = pptx_to_images(file)
        print(file)
        li = upload_to_s3_and_save_to_db(images)
        return li
    else:
        return 'Allowed file types are .pptx'


if __name__ == "__main__":
    app.run(debug=True)
